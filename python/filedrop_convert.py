"""Convert a file to markdown via markitdown's LLM image-description path.

Runs against an OpenAI-compatible gateway. Configuration is read from the
environment; the file path is the first positional argument. The plugin invokes
this as `python -c <inlined source> <path>`, so argv[1] is the file path.
"""

import os
import re
import sys

from markitdown import MarkItDown

# Imported lazily in _make_client so the no-LLM conversion path (plain markitdown
# text extraction) doesn't pay the OpenAI SDK import cost. Tests patch this symbol,
# so it must exist at module scope; the lazy import only fills it when actually
# needed and the attribute is still absent.
OpenAI = None

# Reasoning ("thinking") models emit their chain-of-thought inline in the
# assistant message before the real answer. We strip it so image descriptions
# stay clean; for non-thinking models these patterns never match and the
# content passes through untouched.
_THINK_BLOCK = re.compile(r"<(think|thinking|reasoning)>.*?</\1>", re.DOTALL | re.IGNORECASE)
_DANGLING_CLOSERS = ("</think>", "</thinking>", "</reasoning>")


def strip_thinking(text):
    """Remove reasoning content from an LLM reply, leaving only the answer."""
    if not text:
        return text
    cleaned = _THINK_BLOCK.sub("", text)
    # Some gateways inject the opening tag via the chat template, so the content
    # starts mid-reasoning and only the closing tag survives. Keep what follows
    # the last closer in that case.
    lowered = cleaned.lower()
    for closer in _DANGLING_CLOSERS:
        idx = lowered.rfind(closer)
        if idx != -1:
            after = cleaned[idx + len(closer):]
            if after.strip():
                cleaned = after        # closer in middle → keep what's after
            else:
                cleaned = cleaned[:idx]  # closer at end → strip it, keep what's before
            break
    return cleaned.strip()


def _install_thinking_filter(client):
    """Wrap the client so chat completions have thinking stripped in place."""
    completions = client.chat.completions
    original_create = completions.create

    def create(*args, **kwargs):
        response = original_create(*args, **kwargs)
        for choice in getattr(response, "choices", None) or []:
            message = getattr(choice, "message", None)
            content = getattr(message, "content", None)
            if message is not None and content:
                message.content = strip_thinking(content)
        return response

    completions.create = create
    return client


_DEFAULT_LLM_TIMEOUT_S = 150.0
_DEFAULT_CONNECT_TIMEOUT_S = 15.0


def _llm_timeout(env):
    """Build an httpx.Timeout from FILEDROP_LLM_TIMEOUT (seconds).

    The OpenAI SDK accepts a flat number, but that applies to *every* phase
    (connect/read/write/pool) at once — a connection that establishes but never
    delivers bytes waits the full budget on the read. Splitting connect out
    means a wedged gateway fails fast instead of hanging for minutes. Falls
    back to a flat float if httpx is unavailable (which only happens in tests
    that stub the openai package).
    """
    raw = env.get("FILEDROP_LLM_TIMEOUT")
    try:
        total = float(raw) if raw else _DEFAULT_LLM_TIMEOUT_S
    except (TypeError, ValueError):
        total = _DEFAULT_LLM_TIMEOUT_S
    try:
        import httpx
    except ImportError:
        return total
    return httpx.Timeout(total, connect=_DEFAULT_CONNECT_TIMEOUT_S)


def _make_client(env, **kwargs):
    # x-api-key is sent alongside the SDK's automatic Authorization: Bearer header
    # because the Siemens gateway requires it; other providers ignore it.
    # max_retries=0 because the SDK's default 2 retries silently triple the
    # wait when a gateway hangs — the user is better served by a fast clean
    # error and the existing TS-side error handling.
    global OpenAI
    if OpenAI is None:
        from openai import OpenAI as _OpenAI
        OpenAI = _OpenAI
    key = env["FILEDROP_LLM_KEY"]
    kwargs.setdefault("timeout", _llm_timeout(env))
    kwargs.setdefault("max_retries", 0)
    return OpenAI(
        api_key=key,
        base_url=env.get("FILEDROP_LLM_URL") or None,
        default_headers={"x-api-key": key, "X-Api-Key": key},
        **kwargs,
    )


def _token_kwargs(env, limit):
    """Build the token-limit kwarg for chat.completions.create, honoring the
    capability detected on the TS side (FILEDROP_LLM_TOKEN_PARAM). 'none' omits
    the limit entirely so models that reject both max_tokens and
    max_completion_tokens still work."""
    param = (env.get("FILEDROP_LLM_TOKEN_PARAM") or "max_tokens").strip()
    if param == "none":
        return {}
    if param not in ("max_tokens", "max_completion_tokens"):
        param = "max_tokens"
    return {param: limit}


def _vision_enabled(env):
    """Whether the configured model accepts image input. Defaults to True when
    unset; the TS settings 'Check' sets FILEDROP_LLM_VISION=0 for text-only models."""
    return (env.get("FILEDROP_LLM_VISION") or "1").strip().lower() not in ("0", "false", "no", "off")


def _temperature_kwargs(env):
    """Return {"temperature": 0} when the model supports it, or {} to omit it.
    FILEDROP_LLM_TEMPERATURE=0 means not supported (e.g. reasoning models)."""
    supported = (env.get("FILEDROP_LLM_TEMPERATURE") or "1").strip().lower() not in ("0", "false", "no", "off")
    return {"temperature": 0} if supported else {}


# Images are sent to the gateway as a base64 data URI. Raw high-resolution
# photos can exceed the gateway's request-size limit and come back as HTTP 413
# ("Request Entity Too Large"). We re-encode large images as a compressed JPEG
# to shrink the payload while keeping full resolution; only if compression alone
# is not enough do we fall back to reducing the dimensions.
_DEFAULT_IMAGE_MAX_BYTES = 4 * 1024 * 1024   # re-encode trigger + payload cap
_DEFAULT_IMAGE_JPEG_QUALITY = 85
_DEFAULT_IMAGE_MIN_DIM = 512                  # don't shrink the longest side below this


def _image_limits(env):
    """Read the image-payload limits, honoring optional env overrides."""
    def _int(name, default):
        try:
            return int(env.get(name) or default)
        except (TypeError, ValueError):
            return default
    return (
        _int("FILEDROP_IMAGE_MAX_BYTES", _DEFAULT_IMAGE_MAX_BYTES),
        _int("FILEDROP_IMAGE_JPEG_QUALITY", _DEFAULT_IMAGE_JPEG_QUALITY),
        _int("FILEDROP_IMAGE_MIN_DIM", _DEFAULT_IMAGE_MIN_DIM),
    )


def _normalize_for_jpeg(pix):
    """Return a pixmap that JPEG can encode: JPEG has no alpha channel and only
    supports gray/RGB, so drop alpha and convert CMYK / other colorspaces to RGB."""
    import fitz  # pymupdf  # type: ignore

    if pix.alpha:
        pix = fitz.Pixmap(pix, 0)  # drop alpha
    if pix.colorspace is None or pix.colorspace.n > 3:  # e.g. CMYK -> RGB
        pix = fitz.Pixmap(fitz.csRGB, pix)
    return pix


def _jpeg_bytes_under_cap(pix, quality, max_bytes, min_dim):
    """Encode a pixmap to compressed JPEG, keeping full resolution when possible.
    If the JPEG is still larger than max_bytes, progressively shrink the
    dimensions (~0.8x per step, never below min_dim on the longest side) and
    re-encode until it fits."""
    import fitz  # pymupdf  # type: ignore

    pix = _normalize_for_jpeg(pix)
    data = pix.tobytes("jpg", jpg_quality=quality)
    while len(data) > max_bytes and max(pix.width, pix.height) > min_dim:
        tw, th = int(pix.width * 0.8), int(pix.height * 0.8)
        # Clamp so the longest side never drops below min_dim.
        longest = max(tw, th)
        if longest < min_dim:
            f = min_dim / max(pix.width, pix.height)
            tw, th = int(pix.width * f), int(pix.height * f)
        pix = fitz.Pixmap(pix, max(1, tw), max(1, th), None)
        data = pix.tobytes("jpg", jpg_quality=quality)
    return data


def _compressed_image_path(path, env):
    """If `path` is a large image, write a compressed-JPEG copy to a temp file and
    return its path so the oversized raw bytes never reach the gateway (avoids
    HTTP 413). Returns None when the image is already small enough, or on any
    failure, so the caller falls back to the original file untouched."""
    max_bytes, quality, min_dim = _image_limits(env)
    try:
        if os.path.getsize(path) <= max_bytes:
            return None
        import fitz  # pymupdf  # type: ignore
        import tempfile

        data = _jpeg_bytes_under_cap(fitz.Pixmap(path), quality, max_bytes, min_dim)
        fd, tmp = tempfile.mkstemp(suffix=".jpg")
        try:
            os.write(fd, data)
        finally:
            os.close(fd)
        return tmp
    except Exception as exc:
        print(f"[filedrop] image compression skipped: {type(exc).__name__}: {exc}", file=sys.stderr)
        return None


def _install_progress_counter(client, total):
    """Wrap the client so every chat-completion call (one per slide/image that
    markitdown's built-in converters describe) reports `[filedrop:page-progress]`
    on stderr. Only installed when `total` is large enough to be worth showing."""
    import threading

    completions = client.chat.completions
    original_create = completions.create
    lock = threading.Lock()
    count = 0

    def create(*args, **kwargs):
        nonlocal count
        response = original_create(*args, **kwargs)
        with lock:
            count += 1
            _emit_progress(count, total)
        return response

    completions.create = create
    return client


# Below this many pages/slides a progress indicator is more noise than signal.
PAGE_PROGRESS_THRESHOLD = 3


def _pptx_slide_count(path):
    """Best-effort slide count for a .pptx, used to decide whether to show a
    page-progress indicator. Returns None on any failure (e.g. python-pptx
    missing or the file is malformed) so callers can skip progress reporting."""
    try:
        from pptx import Presentation  # already a markitdown dependency
        return len(Presentation(path).slides)
    except Exception:
        return None


def build_converter(env, progress_total=None):
    client = _make_client(env)
    _install_thinking_filter(client)
    if progress_total and progress_total > PAGE_PROGRESS_THRESHOLD:
        _install_progress_counter(client, progress_total)
    kwargs = {"llm_client": client, "llm_model": env["FILEDROP_LLM_MODEL"]}
    prompt = env.get("FILEDROP_LLM_PROMPT")
    if prompt:
        kwargs["llm_prompt"] = prompt
    return MarkItDown(**kwargs)


def _emit_phase(phase):
    print(f"[filedrop:phase] {phase}", file=sys.stderr, flush=True)


def _emit_progress(current, total):
    print(f"[filedrop:page-progress] {current}/{total}", file=sys.stderr, flush=True)


def _llm_configured(env):
    """Whether an LLM gateway is configured. The TS side only runs this script
    with these env vars set (or for the describe path), but guarding here lets
    convert() still extract a file's text via plain markitdown when they are
    absent — and skips importing the OpenAI SDK entirely on that fast path."""
    return bool(env.get("FILEDROP_LLM_KEY") and env.get("FILEDROP_LLM_MODEL"))


def _with_llm_error(exc, plain_text):
    """Prepend an error callout to the plain markitdown fallback text."""
    callout = f"> [!error] LLM conversion failed — plain text only\n> {type(exc).__name__}: {exc}"
    return f"{callout}\n\n{plain_text}" if plain_text and plain_text.strip() else callout


def convert(path, env):
    # markitdown hands the path to puremagic, which raises an uncaught
    # PureError("Not a regular file") for anything that isn't a regular file
    # (e.g. a directory), aborting the whole conversion with a cryptic message.
    # Guard up front so the failure is legible instead.
    if not os.path.isfile(path):
        kind = "a directory" if os.path.isdir(path) else "missing"
        print(f"[filedrop] input path is {kind}, not a regular file: {path}", file=sys.stderr)
        return ""

    is_pdf = path.lower().endswith(".pdf")

    # No gateway configured: plain markitdown text extraction. Keeps trivial
    # conversions fast (the OpenAI SDK is never imported) and lets convert() work
    # without an LLM at all.
    if not _llm_configured(env):
        _emit_phase("markitdown")
        return _convert_without_llm(path)

    # Try markitdown with LLM support first. build_converter passes llm_client
    # and llm_model to MarkItDown, so embedded images inside the PDF are
    # described by the LLM. Works well for text-layer PDFs. Scanned PDFs often
    # return empty here because pdfminer finds no text layer and there are no
    # discrete embedded images — we catch that and fall through.
    # Pure-image files go straight into the LLM via markitdown's image path,
    # so signal that phase up front instead of "markitdown".
    image_exts = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif")
    is_image = path.lower().endswith(image_exts)
    _emit_phase("llm-image" if is_image else "markitdown")

    # For large images, hand markitdown a compressed-JPEG copy so the oversized
    # raw bytes never reach the gateway (avoids HTTP 413). Small images and other
    # file types use the original path untouched.
    convert_path = path
    temp_image = _compressed_image_path(path, env) if is_image else None
    if temp_image:
        convert_path = temp_image
    progress_total = _pptx_slide_count(path) if path.lower().endswith(".pptx") else None
    try:
        result = build_converter(env, progress_total=progress_total).convert(convert_path).text_content
        if result and result.strip():
            return result
    except Exception as exc:
        url = env.get("FILEDROP_LLM_URL", "(not set)")
        model = env.get("FILEDROP_LLM_MODEL", "(not set)")
        print(f"[filedrop] markitdown LLM step failed [{model} @ {url}]: {type(exc).__name__}: {exc}", file=sys.stderr)
        if not is_pdf:
            plain = _convert_without_llm(convert_path)
            return _with_llm_error(exc, plain)
    finally:
        if temp_image:
            try:
                os.remove(temp_image)
            except OSError:
                pass

    # For PDFs that produced nothing (scanned / image-only), render each page
    # via PyMuPDF and OCR with the LLM.
    if is_pdf:
        _emit_phase("llm-image")
        try:
            result = _convert_pdf_pages_with_llm(path, env)
        except Exception as exc:
            plain = _convert_without_llm(path)
            return _with_llm_error(exc, plain)
        if result:
            return result

    return ""


def _convert_without_llm(path):
    """Re-run markitdown with no LLM client so a file's extractable text is still
    captured when the LLM-enhanced step fails. Returns the plain text, or "" if
    even this fails — never raises, so a failure here can't crash the process and
    dump the whole invocation into the note."""
    try:
        result = MarkItDown().convert(path).text_content
        if result and result.strip():
            return result
    except Exception as exc:
        print(f"[filedrop] plain markitdown fallback failed: {type(exc).__name__}: {exc}", file=sys.stderr)
    return ""


def _convert_pdf_pages_with_llm(path, env):
    """Render every PDF page via PyMuPDF and OCR with the LLM.

    Returns combined markdown, or empty string if PyMuPDF is unavailable
    so the caller falls back to plain markitdown text extraction.
    """
    import base64
    import threading
    from concurrent.futures import ThreadPoolExecutor

    if not _vision_enabled(env):
        return ("> [!warning] Scanned-PDF OCR skipped — the configured model has no "
                "vision (image input) support.")

    try:
        import fitz  # pymupdf  # type: ignore
    except ImportError:
        print("[filedrop] PyMuPDF not installed — cannot OCR scanned PDF pages", file=sys.stderr)
        return ""

    client = _make_client(env)
    _install_thinking_filter(client)

    prompt = (
        env.get("FILEDROP_LLM_PROMPT")
        or "Transcribe all text visible on this page exactly as it appears. "
           "For any charts, tables, or images also provide a concise description."
    )

    try:
        doc = fitz.open(path)
    except Exception:
        return ""

    if not doc.page_count:
        return ""

    # Render every page up front: PyMuPDF documents are not safe to access from
    # multiple threads, and rendering is fast/CPU-bound. The slow part is the
    # per-page LLM request, which we then run in a small thread pool. The OpenAI
    # client (httpx under the hood) is safe to share across threads.
    # Each page is a compressed JPEG instead of raw PNG: far smaller payload
    # (avoids HTTP 413 on large/high-res pages), with a dimension fallback if
    # still big.
    max_bytes, quality, min_dim = _image_limits(env)
    mat = fitz.Matrix(2, 2)  # 2× zoom for better OCR quality
    images = [
        base64.b64encode(
            _jpeg_bytes_under_cap(page.get_pixmap(matrix=mat), quality, max_bytes, min_dim)
        ).decode("ascii")
        for page in doc
    ]

    total = len(images)
    show_progress = total > PAGE_PROGRESS_THRESHOLD
    progress_lock = threading.Lock()
    completed = 0
    stop_event = threading.Event()

    def _ocr_page(item):
        nonlocal completed
        page_num, img_b64 = item
        if stop_event.is_set():
            return None
        try:
            resp = client.chat.completions.create(
                model=env["FILEDROP_LLM_MODEL"],
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                    ],
                }],
                **_token_kwargs(env, 4096),
                **_temperature_kwargs(env),
            )
            text = (resp.choices[0].message.content or "").strip()
            text = strip_thinking(text)
        except Exception as exc:
            stop_event.set()
            print(f"[filedrop] PDF OCR failed on page {page_num}: {type(exc).__name__}: {exc}", file=sys.stderr)
            raise
        if show_progress:
            with progress_lock:
                completed += 1
                _emit_progress(completed, total)
        return f"### Page {page_num}\n\n{text}"

    items = list(enumerate(images, 1))
    max_workers = min(4, len(items))
    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        # pool.map preserves input order, so pages stay in page order regardless
        # of which request finishes first; progress is still reported in
        # completion order since that reflects actual work done.
        pages = list(pool.map(_ocr_page, items))

    return "\n\n".join(p for p in pages if p is not None)


DEFAULT_DESCRIBE_PROMPT = (
    "A user saved a file named '{filename}' that cannot be converted to text. "
    "Based only on its filename, briefly explain what this file most likely is "
    "(for example, an installer for a specific application, or a system/standard "
    "library). Keep it to 1-3 sentences and make clear it is an educated guess."
)


def describe(path, env):
    """Ask the LLM what a file likely is, given only its filename."""
    client = _make_client(env)
    _install_thinking_filter(client)
    filename = os.path.basename(path)
    prompt = (env.get("FILEDROP_DESCRIBE_PROMPT") or DEFAULT_DESCRIBE_PROMPT).format(filename=filename)
    response = client.chat.completions.create(
        model=env["FILEDROP_LLM_MODEL"],
        messages=[{"role": "user", "content": prompt}],
        **_token_kwargs(env, 2048),
        **_temperature_kwargs(env),
    )
    return response.choices[0].message.content or ""


# ---------------------------------------------------------------------------
# Structure-aware PPTX extraction
#
# markitdown flattens a slide's shapes into plain paragraphs in stored z-order
# and emits `![desc](PictureN.jpg)` references whose image bytes it never writes.
# Here we read the deck with python-pptx instead: every shape keeps its geometry
# (so the TS side can reconstruct columns / reading order via the LLM) and every
# embedded picture's bytes are written to a temp dir under the *same* name
# markitdown would reference, so the links resolve once copied into the vault.
# ---------------------------------------------------------------------------

DEFAULT_IMAGE_PROMPT = (
    "Describe this image concisely for use as markdown alt text. "
    "One or two sentences focused on the meaningful content; no preamble."
)


def _iter_shapes(shapes, inherited_geom=None):
    """Yield (shape, effective_geom) for every shape, recursing into group shapes
    (depth-first, document order) so pictures/text nested inside groups aren't
    missed. A grouped child's own left/top are in the group's *child* coordinate
    space, not slide EMUs, so they can't be compared with top-level shapes; we
    attribute the (outermost) group's slide-space geometry to all its descendants
    instead, which keeps a group clustered in its real slide region. MSO_SHAPE_TYPE
    .GROUP == 6; checking the attribute as well stays robust across versions."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            group_geom = inherited_geom if inherited_geom is not None else _shape_geometry(shape)
            yield from _iter_shapes(shape.shapes, group_geom)
        else:
            yield shape, inherited_geom if inherited_geom is not None else _shape_geometry(shape)


def _emu(value):
    """Coerce a python-pptx length (EMU) to int, or None when unset."""
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _shape_geometry(shape):
    return {
        "left": _emu(getattr(shape, "left", None)),
        "top": _emu(getattr(shape, "top", None)),
        "width": _emu(getattr(shape, "width", None)),
        "height": _emu(getattr(shape, "height", None)),
    }


def _geoms_overlap(a, b):
    """Return True if two geom dicts (left/top/width/height in EMU) intersect."""
    if any(v is None for v in (a.get("left"), a.get("top"), a.get("width"), a.get("height"),
                                b.get("left"), b.get("top"), b.get("width"), b.get("height"))):
        return False
    return (a["left"] < b["left"] + b["width"] and a["left"] + a["width"] > b["left"] and
            a["top"]  < b["top"]  + b["height"] and a["top"]  + a["height"] > b["top"])


def _picture_filename(shape):
    """The stripped-name base markitdown uses for a PPTX picture, plus .jpg.
    Not used as the on-disk name directly — see _unique_picture_name, which
    prefixes the slide index to keep names unique across the deck."""
    return re.sub(r"\W", "", shape.name or "image") + ".jpg"


def _unique_picture_name(shape, slide_index, seen):
    """A deck-unique picture filename. PowerPoint shape names ("Picture 3") are
    only unique within a slide, so prefix the slide index; on the rare residual
    clash append a counter. We control the refs end to end (the LLM echoes them,
    rewriteImageLinks matches them), so the name need not match markitdown."""
    base = re.sub(r"\W", "", shape.name or "image")
    stem = f"slide{slide_index}-{base}"
    name = f"{stem}.jpg"
    counter = 2
    while name in seen:
        name = f"{stem}-{counter}.jpg"
        counter += 1
    seen.add(name)
    return name


def _picture_alt_text(shape):
    """The author-provided alt text on a picture, if any (matches markitdown)."""
    try:
        return shape._element._nvXxPr.cNvPr.get("descr") or ""
    except Exception:
        return ""


def _table_markdown(table):
    """Render a pptx table as a GitHub-flavored markdown table."""
    rows = []
    for row in table.rows:
        rows.append([" ".join((cell.text or "").split()) for cell in row.cells])
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    lines = ["| " + " | ".join(rows[0]) + " |",
             "| " + " | ".join("---" for _ in range(width)) + " |"]
    for r in rows[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _chart_title(chart):
    """A chart's title text, or "" — best-effort (never raises)."""
    try:
        if chart.has_title:
            return (chart.chart_title.text_frame.text or "").strip()
    except Exception:
        pass
    return ""


def _chart_markdown(chart):
    """Render a chart's underlying data as a markdown table (categories × series)
    so charts aren't silently dropped — the reflow LLM can table it and describe
    the trend. Returns "" when no series data is available."""
    def _fmt(v):
        if v is None:
            return ""
        if isinstance(v, float) and v.is_integer():
            return str(int(v))
        return str(v)

    try:
        categories = list(chart.plots[0].categories) if chart.plots else []
    except Exception:
        categories = []
    series = []
    try:
        for idx, s in enumerate(chart.series, 1):
            series.append((str(s.name) if s.name else f"Series {idx}", list(s.values)))
    except Exception:
        series = []
    if not series:
        return ""

    header = ["Category"] + [name for name, _ in series]
    lines = ["| " + " | ".join(header) + " |",
             "| " + " | ".join("---" for _ in header) + " |"]
    row_count = max((len(vals) for _, vals in series), default=0)
    for r in range(row_count):
        cat = str(categories[r]) if r < len(categories) else ""
        row = [cat] + [_fmt(vals[r]) if r < len(vals) else "" for _, vals in series]
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _slide_notes(slide):
    """The slide's speaker-notes text, or "" — markitdown surfaces these, so the
    structured path must too to avoid dropping content."""
    try:
        if not slide.has_notes_slide:
            return ""
        return (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        return ""


def _describe_image(image, env, client, prompt, overlay_texts=None):
    """Return a short LLM description of a python-pptx Image, or "" on failure.
    Large blobs are re-encoded as compressed JPEG (reusing the PDF-page 413
    guard) so the payload stays under the gateway's request-size limit.
    overlay_texts: text labels from PowerPoint shapes positioned over this image."""
    import base64

    blob = image.blob
    content_type = getattr(image, "content_type", None) or "image/png"
    max_bytes, quality, min_dim = _image_limits(env)
    if len(blob) > max_bytes:
        try:
            import fitz  # pymupdf  # type: ignore
            import tempfile

            fd, tmp = tempfile.mkstemp()
            try:
                os.write(fd, blob)
            finally:
                os.close(fd)
            try:
                blob = _jpeg_bytes_under_cap(fitz.Pixmap(tmp), quality, max_bytes, min_dim)
                content_type = "image/jpeg"
            finally:
                os.remove(tmp)
        except Exception as exc:
            print(f"[filedrop] pptx image compression skipped: {type(exc).__name__}: {exc}", file=sys.stderr)

    b64 = base64.b64encode(blob).decode("ascii")
    user_content = [
        {"type": "text", "text": prompt},
        {"type": "image_url", "image_url": {"url": f"data:{content_type};base64,{b64}"}},
    ]
    if overlay_texts:
        label_str = ", ".join(f'"{t}"' for t in overlay_texts)
        user_content.append({
            "type": "text",
            "text": (
                f"Note: the following text labels appear as separate PowerPoint shapes "
                f"positioned over this image (they are not burned into the image pixels): "
                f"{label_str}. Incorporate these labels into your description."
            ),
        })
    try:
        resp = client.chat.completions.create(
            model=env["FILEDROP_LLM_MODEL"],
            messages=[{"role": "user", "content": user_content}],
            **_token_kwargs(env, 1024),
            **_temperature_kwargs(env),
        )
        return strip_thinking((resp.choices[0].message.content or "").strip())
    except Exception as exc:
        print(f"[filedrop] pptx image description failed: {type(exc).__name__}: {exc}", file=sys.stderr)
        return ""


def _extract_slide_elements(slide, pics_dir, slide_index, seen):
    """Walk one slide's shapes (recursing groups) into ordered elements with
    geometry, writing each picture's bytes to pics_dir under a deck-unique name
    (`seen` tracks names already used across the whole deck). Returns
    (elements, pending_descriptions) where pending is a list of
    (element, pptx_image) whose description still needs an LLM call."""
    elements = []
    pending = []
    for shape, geom in _iter_shapes(slide.shapes):
        image = None
        try:
            image = shape.image  # raises on non-picture shapes
        except Exception:
            image = None
        if image is not None:
            filename = _unique_picture_name(shape, slide_index, seen)
            dest = os.path.join(pics_dir, filename)
            try:
                with open(dest, "wb") as fh:
                    fh.write(image.blob)
            except OSError as exc:
                print(f"[filedrop] failed writing pptx image {filename}: {exc}", file=sys.stderr)
                continue
            alt = _picture_alt_text(shape)
            element = {"type": "picture", "geom": geom, "filename": filename, "description": alt}
            elements.append(element)
            if not alt:
                pending.append((element, image))
            continue

        if getattr(shape, "has_table", False):
            try:
                md = _table_markdown(shape.table)
            except Exception as exc:
                print(f"[filedrop] pptx table access failed on slide {slide_index}: {type(exc).__name__}: {exc}", file=sys.stderr)
                continue
            if md:
                elements.append({"type": "table", "geom": geom, "markdown": md})
            continue

        if getattr(shape, "has_chart", False):
            try:
                chart = shape.chart
                title = _chart_title(chart)
                md = _chart_markdown(chart)
            except Exception as exc:
                print(f"[filedrop] pptx chart access failed on slide {slide_index}: {type(exc).__name__}: {exc}", file=sys.stderr)
                continue
            if title or md:
                elements.append({"type": "chart", "geom": geom, "title": title, "markdown": md})
            continue

        if getattr(shape, "has_text_frame", False):
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs) if para.runs else (para.text or "")
                text = text.strip()
                if text:
                    paragraphs.append({"text": text, "level": getattr(para, "level", 0) or 0})
            if paragraphs:
                elements.append({"type": "text", "geom": geom, "paragraphs": paragraphs})
            continue

        # Structural drawing shapes (arrows, connectors, bare boxes, etc.) that
        # carry no text frame — capture them so the reflow LLM can describe
        # the diagram structure they form.
        shape_type_name = getattr(getattr(shape, "shape_type", None), "name", None)
        if shape_type_name in ("AUTO_SHAPE", "LINE", "CONNECTOR", "FREEFORM"):
            kind = getattr(getattr(shape, "auto_shape_type", None), "name", shape_type_name)
            elements.append({"type": "shape", "geom": geom, "shape_kind": kind.lower(), "text": ""})

    # Post-processing: for each image awaiting LLM description, collect text
    # elements that spatially overlap it (overlay labels whose text is a
    # separate PowerPoint shape positioned over the image). Pass them to the
    # vision LLM as context, but keep them in the element list so they remain
    # visible in the output.
    overlay_map = {}
    for pi, (pic_el, img_obj) in enumerate(pending):
        texts = []
        for el in elements:
            if el["type"] != "text":
                continue
            if _geoms_overlap(pic_el["geom"], el["geom"]):
                for para in el.get("paragraphs", []):
                    t = para.get("text", "").strip()
                    if t:
                        texts.append(t)
        overlay_map[pi] = texts

    pending = [(el, img, overlay_map[pi]) for pi, (el, img) in enumerate(pending)]

    return elements, pending


def convert_pptx_structured(path, env):
    """Extract a .pptx into structured per-slide JSON plus its embedded images.

    Returns {"slides": [...], "pictures": [{"filename", "path"}]} where pictures
    live in a fresh temp dir (the TS side copies them into the vault and removes
    the dir). Returns None to signal the caller should fall back to markitdown
    (python-pptx unavailable, or the deck can't be opened)."""
    try:
        from pptx import Presentation
    except ImportError:
        print("[filedrop] python-pptx not installed — falling back to markitdown for pptx", file=sys.stderr)
        return None
    try:
        prs = Presentation(path)
    except Exception as exc:
        print(f"[filedrop] could not open pptx: {type(exc).__name__}: {exc}", file=sys.stderr)
        return None

    import tempfile

    pics_dir = tempfile.mkdtemp(prefix="filedrop_pptx_")
    describe_enabled = _llm_configured(env) and _vision_enabled(env)
    client = None
    if describe_enabled:
        try:
            client = _make_client(env)
            _install_thinking_filter(client)
        except Exception as exc:
            print(f"[filedrop] LLM client unavailable for pptx images: {type(exc).__name__}: {exc}", file=sys.stderr)
            describe_enabled = False

    image_prompt = env.get("FILEDROP_LLM_PROMPT") or DEFAULT_IMAGE_PROMPT

    _emit_phase("markitdown")
    out_slides = []
    pictures = []
    pending = []
    seen_names = set()
    for idx, slide in enumerate(prs.slides, 1):
        try:
            elements, slide_pending = _extract_slide_elements(slide, pics_dir, idx, seen_names)
        except Exception as exc:
            print(f"[filedrop] pptx slide {idx} extraction failed: {type(exc).__name__}: {exc}", file=sys.stderr)
            elements, slide_pending = [], []
        for element in elements:
            if element.get("type") == "picture":
                pictures.append({"filename": element["filename"], "path": os.path.join(pics_dir, element["filename"])})
        slide_doc = {"index": idx, "elements": elements}
        notes = _slide_notes(slide)
        if notes:
            slide_doc["notes"] = notes
        out_slides.append(slide_doc)
        pending.extend(slide_pending)

    if describe_enabled and client is not None and pending:
        import threading
        from concurrent.futures import ThreadPoolExecutor

        _emit_phase("llm-image")
        total = len(pending)
        lock = threading.Lock()
        completed = 0
        _emit_progress(0, total)

        def _run(item):
            nonlocal completed
            element, image, overlay_texts = item
            element["description"] = _describe_image(image, env, client, image_prompt, overlay_texts=overlay_texts)
            with lock:
                completed += 1
                _emit_progress(completed, total)

        with ThreadPoolExecutor(max_workers=min(4, total)) as pool:
            list(pool.map(_run, pending))

    return {"slides": out_slides, "pictures": pictures}


def main(argv=None, env=None):
    argv = sys.argv if argv is None else argv
    env = os.environ if env is None else env
    if len(argv) > 2 and argv[1] == "--pptx":
        import json

        result = convert_pptx_structured(argv[2], env)
        # Empty stdout signals the TS side to fall back to markitdown.
        if result is not None:
            sys.stdout.buffer.write(json.dumps(result).encode("utf-8"))
        return
    path = argv[1]
    if env.get("FILEDROP_DESCRIBE"):
        result = describe(path, env)
    else:
        ext = os.path.splitext(path)[1].lower()
        describe_exts = {e.strip().lower() for e in env.get("FILEDROP_DESCRIBE_EXTS", "").split(",") if e.strip()}
        if ext and ext in describe_exts:
            desc = describe(path, env)
            result = (
                "> [!warning] Unsupported file format — could not convert\n"
                "> markitdown can't convert this file type. Best guess from the filename (may be inaccurate):\n"
                "\n"
                + desc
            ) if desc else ""
        else:
            result = convert(path, env)
    sys.stdout.buffer.write(result.encode("utf-8"))


if __name__ == "__main__":
    main()
